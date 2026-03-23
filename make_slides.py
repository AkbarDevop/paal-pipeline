#!/usr/bin/env python3
"""Generate weekly meeting PowerPoint slides for PAAL posture classification."""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from config import OUTPUT_DIR

# ── Colors ───────────────────────────────────────────────────────────────────

WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT     = RGBColor(0x16, 0x87, 0xA7)
ACCENT2    = RGBColor(0x11, 0x99, 0x8E)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
MID_GRAY   = RGBColor(0x66, 0x66, 0x77)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE     = RGBColor(0xF3, 0x9C, 0x12)
BLUE       = RGBColor(0x29, 0x80, 0xB9)


# ── Helpers ──────────────────────────────────────────────────────────────────

def img_path(name):
    """Get path to an output image, return None if missing."""
    p = os.path.join(OUTPUT_DIR, name)
    return p if os.path.exists(p) else None


def add_bg(slide, color=DARK):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_image(slide, path, left, top, width=None, height=None):
    """Add an image to a slide. Returns the shape or None if file missing."""
    if not path or not os.path.exists(path):
        return None
    kwargs = {}
    if width:
        kwargs["width"] = Inches(width)
    if height:
        kwargs["height"] = Inches(height)
    return slide.shapes.add_picture(path, Inches(left), Inches(top), **kwargs)


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, size=16,
                    color=WHITE, spacing=Pt(8)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_table(slide, left, top, width, height, rows, cols):
    """Add a table shape."""
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    return table_shape.table


def style_cell(cell, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.CENTER,
               fill_color=None):
    """Style a table cell."""
    cell.text = str(text)
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def slide_title(slide, title, subtitle=None):
    """Add title bar to a dark slide."""
    add_text(slide, 0.6, 0.3, 8.8, 0.7, title, size=32, color=WHITE, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(1.0), Inches(2.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    if subtitle:
        add_text(slide, 0.6, 1.1, 8.8, 0.4, subtitle, size=16, color=MID_GRAY)


# ── Slides ───────────────────────────────────────────────────────────────────

def make_title_slide(prs):
    """Slide 1: Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK)

    add_text(slide, 0.6, 1.5, 8.8, 1.0,
             "PAAL Sow Posture Classification",
             size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, 0.6, 2.4, 8.8, 0.6,
             "OAK-D ToF Camera  |  Deep Learning Pipeline",
             size=20, color=ACCENT, align=PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(3.5), Inches(3.2), Inches(3.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_text(slide, 0.6, 3.5, 8.8, 0.5,
             "Weekly Progress Report — Task A",
             size=18, color=MID_GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, 0.6, 4.2, 8.8, 0.4,
             "Akbar",
             size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)


def make_pipeline_slide(prs):
    """Slide 2: Pipeline overview — with pipeline diagram image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    slide_title(slide, "Pipeline Overview", "End-to-end automated pipeline")

    # Use the generated pipeline diagram
    p = img_path("pipeline_overview.png")
    if p:
        add_image(slide, p, 0.3, 1.5, width=9.4)
    else:
        # Fallback to text boxes
        steps = [
            ("1", "Scan Data", "Build metadata from\nOAK-D folders", ACCENT),
            ("2", "Prefilter", "Depth-based pig\npresence detection", ACCENT),
            ("3", "Label", "Manual posture\nlabeling (GUI)", ACCENT2),
            ("4", "Crop", "Remove neighbors\n& ceiling", ACCENT),
            ("5", "Train", "3 backbones ×\n3 modalities", GREEN),
            ("6", "Evaluate", "Test accuracy &\nconfusion matrices", BLUE),
            ("7", "Visualize", "Summary plots\nfor presentation", ORANGE),
        ]
        x_start, box_w, gap, y_top = 0.3, 1.2, 0.1, 1.8
        for i, (num, name, desc, color) in enumerate(steps):
            x = x_start + i * (box_w + gap)
            add_rect(slide, x, y_top, box_w, 2.2, RGBColor(0x25, 0x25, 0x3A), color)
            add_text(slide, x, y_top + 0.1, box_w, 0.4, num, size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, x, y_top + 0.5, box_w, 0.4, name, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, x + 0.05, y_top + 1.0, box_w - 0.1, 1.0, desc, size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, 0.6, 4.6, 8.8, 0.5,
             "python run_pipeline.py --all    # runs the full pipeline automatically",
             size=14, color=ACCENT, align=PP_ALIGN.CENTER, font_name="Consolas")


def make_posture_examples(prs):
    """Posture examples — actual RGB photos of standing/sitting/lying."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    slide_title(slide, "Three Posture Classes", "What the model learns to distinguish")

    p = img_path("posture_examples_rgb.png")
    if p:
        add_image(slide, p, 0.3, 1.3, width=9.4)

    add_bullet_list(slide, 0.6, 4.2, 8.8, 1.2, [
        "Standing: pig upright on all 4 legs — visible leg gap under body",
        "Sitting: pig on haunches with front legs extended — rarest posture (3.8% of data)",
        "Lying: pig on side or belly — most common posture, especially at night",
    ], size=13, color=LIGHT_GRAY)


def make_modality_slide(prs):
    """3 modalities × 3 postures grid — cropped images."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    slide_title(slide, "3 Modalities × 3 Postures",
                "RGB, IR, and Depth — each captures different information")

    p = img_path("posture_modality_grid.png")
    if p:
        add_image(slide, p, 2.0, 1.2, width=6.0)

    add_bullet_list(slide, 0.3, 4.6, 9.4, 0.8, [
        "RGB: color + texture  |  IR: works in darkness (thermal)  |  Depth: 3D shape (distance to camera)",
    ], size=12, color=MID_GRAY)


def make_crop_slide(prs):
    """Crop before/after — shows green box on original, then cropped result."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    slide_title(slide, "Image Cropping", "Remove neighbor pigs and ceiling hardware")

    p = img_path("crop_before_after.png")
    if p:
        add_image(slide, p, 0.3, 1.3, width=9.4)

    add_bullet_list(slide, 0.6, 4.5, 8.8, 1.0, [
        "Green box: fixed crop region based on stall bar positions (same for all frames)",
        "Cropped to 224×224 input — model sees only the target pig, not distractions",
    ], size=12, color=LIGHT_GRAY)


def make_data_slide(prs):
    """Slide 4: Dataset overview — with dataset chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    slide_title(slide, "Dataset", "20 pigs × 3 modalities, pig-ID-based split")

    # Left side: table
    tbl = add_table(slide, 0.4, 1.5, 5.0, 2.2, 5, 6)

    headers = ["Split", "Pig IDs", "Samples", "Standing", "Sitting", "Lying"]
    data = [
        ["Train",  "0–11",   "918",   "451", "40",  "427"],
        ["Val",    "12–15",  "264",   "121", "8",   "135"],
        ["Test",   "16–19",  "190",   "84",  "4",   "102"],
        ["Total",  "20 pigs","1,372", "656", "52",  "664"],
    ]

    for j, h in enumerate(headers):
        style_cell(tbl.cell(0, j), h, size=11, bold=True, color=WHITE, fill_color=ACCENT)
    for i, row in enumerate(data):
        fill = RGBColor(0x25, 0x25, 0x3A) if i % 2 == 0 else RGBColor(0x20, 0x20, 0x35)
        if i == 3:
            fill = RGBColor(0x1C, 0x3D, 0x5A)
        for j, val in enumerate(row):
            style_cell(tbl.cell(i + 1, j), val, size=10, color=WHITE, bold=(i == 3), fill_color=fill)

    # Right side: dataset overview chart
    p = img_path("dataset_overview.png")
    if p:
        add_image(slide, p, 5.5, 1.4, width=4.3)

    add_bullet_list(slide, 0.4, 3.9, 9.2, 1.6, [
        "Pig-ID split: test pigs never seen during training (no data leakage)",
        "Class imbalance: sitting = 3.8% of data — addressed with class-weighted loss",
        "Left: Sowbot 2022 dataset (9,072 IR images) used for pre-training experiment",
    ], size=12, color=LIGHT_GRAY)


def make_models_slide(prs):
    """Slide 5: Model architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    slide_title(slide, "Model Architecture", "3 pretrained backbones × 3 input modalities")

    tbl = add_table(slide, 0.6, 1.5, 8.8, 2.0, 4, 4)

    headers = ["Backbone", "Parameters", "ImageNet Top-1", "Strengths"]
    data = [
        ["MobileNetV2",  "3.4M",   "71.8%", "Lightweight, fast, deployable"],
        ["Xception",     "22.9M",  "79.0%", "Depthwise separable convolutions"],
        ["DenseNet121",  "8.0M",   "74.4%", "Feature reuse, gradient flow"],
    ]

    for j, h in enumerate(headers):
        style_cell(tbl.cell(0, j), h, size=13, bold=True, color=WHITE, fill_color=ACCENT)
    for i, row in enumerate(data):
        fill = RGBColor(0x25, 0x25, 0x3A) if i % 2 == 0 else RGBColor(0x20, 0x20, 0x35)
        for j, val in enumerate(row):
            align = PP_ALIGN.LEFT if j == 3 else PP_ALIGN.CENTER
            style_cell(tbl.cell(i + 1, j), val, size=12, color=WHITE, fill_color=fill, align=align)

    add_bullet_list(slide, 0.6, 3.8, 8.8, 1.6, [
        "Transfer learning: ImageNet pretrained → fine-tuned on OAK-D data",
        "Cropped inputs (224×224): removes neighboring pigs and ceiling hardware",
        "Class-weighted cross-entropy loss to boost sitting class learning",
        "Adam optimizer, ReduceLROnPlateau scheduler, 30 epochs",
    ], size=13, color=LIGHT_GRAY)


def make_results_slide(prs):
    """Slide 6: Results table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    slide_title(slide, "Results: 98.4% Test Accuracy",
                "Held-out test set (pigs 16–19, 190 frames)")

    tbl = add_table(slide, 0.6, 1.5, 4.5, 2.0, 4, 4)

    headers = ["Backbone", "RGB", "IR", "Depth"]
    data = [
        ["MobileNetV2",  "98.4%", "98.4%", "98.4%"],
        ["Xception",     "96.8%", "98.4%", "97.9%"],
        ["DenseNet121",  "98.4%", "97.9%", "98.4%"],
    ]

    for j, h in enumerate(headers):
        style_cell(tbl.cell(0, j), h, size=12, bold=True, color=WHITE, fill_color=ACCENT)
    for i, row in enumerate(data):
        fill = RGBColor(0x25, 0x25, 0x3A) if i % 2 == 0 else RGBColor(0x20, 0x20, 0x35)
        for j, val in enumerate(row):
            c, bold = WHITE, False
            if j > 0 and val == "98.4%":
                c, bold = GREEN, True
            style_cell(tbl.cell(i + 1, j), val, size=12, color=c, bold=bold, fill_color=fill)

    # Right side: comparison bar chart
    p = img_path("modality_comparison_test.png")
    if p:
        add_image(slide, p, 5.3, 1.3, width=4.5)

    add_bullet_list(slide, 0.6, 3.8, 8.8, 1.6, [
        "6 out of 9 models achieve 98.4% (only 3 misclassified out of 190)",
        "MobileNetV2/Depth: best sitting recall (4/4 = 100%) with smallest model",
        "All modalities viable — depth and IR work even in low-light conditions",
    ], size=12, color=LIGHT_GRAY)


def make_confusion_slide(prs):
    """Slide 7: Confusion matrices — actual plots."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    slide_title(slide, "Confusion Matrices — Top Models",
                "98.4% accuracy = only 3 errors out of 190 test frames")

    p = img_path("top_confusion_matrices.png")
    if p:
        add_image(slide, p, 0.3, 1.4, width=9.4)

    add_bullet_list(slide, 0.6, 4.4, 8.8, 1.0, [
        "Main errors: 1 standing→sitting, 1 sitting→lying, 1 lying→standing",
        "MobileNetV2/Depth gets all 4 sitting frames correct (best sitting recall)",
    ], size=12, color=LIGHT_GRAY)


def make_sitting_slide(prs):
    """Slide 8: Sitting class deep dive — with F1 chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    slide_title(slide, "Sitting Class Analysis",
                "Main challenge: only 52 sitting samples (3.8% of data)")

    # Left: table
    tbl = add_table(slide, 0.4, 1.5, 4.8, 2.0, 4, 4)
    headers = ["Model", "Precision", "Recall", "F1"]
    data = [
        ["MobNetV2/Depth", "0.67", "1.00 (4/4)", "0.80"],
        ["MobNetV2/RGB",   "0.75", "0.75 (3/4)", "0.75"],
        ["DenseNet/IR",    "0.67", "0.50 (2/4)", "0.57"],
    ]

    for j, h in enumerate(headers):
        style_cell(tbl.cell(0, j), h, size=11, bold=True, color=WHITE, fill_color=ACCENT)
    for i, row in enumerate(data):
        fill = RGBColor(0x25, 0x25, 0x3A) if i % 2 == 0 else RGBColor(0x20, 0x20, 0x35)
        for j, val in enumerate(row):
            c, bold = WHITE, False
            if i == 0 and j > 0:
                c, bold = GREEN, True
            style_cell(tbl.cell(i + 1, j), val, size=11, color=c, bold=bold, fill_color=fill)

    # Right: sitting F1 comparison chart
    p = img_path("comparison_sitting_f1.png")
    if p:
        add_image(slide, p, 5.3, 1.3, width=4.5)

    add_bullet_list(slide, 0.4, 3.8, 9.2, 1.6, [
        "MobileNetV2/Depth correctly classifies all 4 sitting test frames",
        "Class-weighted loss crucial: without it, model ignores rare sitting class",
        "More sitting labels would further improve reliability (currently 52 total)",
    ], size=12, color=LIGHT_GRAY)


def make_sowbot_slide(prs):
    """Slide 9: Sowbot pre-training — with comparison chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    slide_title(slide, "Experiment: Sowbot Pre-training",
                "Pre-train on 9,072 sowbot IR images → fine-tune on OAK-D")

    # Left: comparison bar chart
    p = img_path("comparison_accuracy.png")
    if p:
        add_image(slide, p, 0.2, 1.4, width=5.5)

    # Right: heatmap
    p2 = img_path("comparison_heatmap.png")
    if p2:
        add_image(slide, p2, 5.5, 1.4, width=4.3)

    add_text(slide, 0.6, 3.8, 8.8, 0.5,
             "Result: Pre-training did NOT improve performance",
             size=18, color=ORANGE, bold=True)

    add_bullet_list(slide, 0.6, 4.3, 8.8, 1.2, [
        "ImageNet features already sufficient for posture classification",
        "Domain gap: sowbot IR (480×480 grayscale) vs OAK-D (640×480 multi-modal)",
        "Ceiling effect: already at 98.4% — little room for improvement",
    ], size=12, color=LIGHT_GRAY)


def make_sowbot_curves_slide(prs):
    """Slide 10: Sowbot pre-training curves."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    slide_title(slide, "Sowbot Pre-training Curves",
                "All 3 backbones reach ~99.9% on sowbot data (9,072 IR images)")

    p = img_path("sowbot_pretrain_curves.png")
    if p:
        add_image(slide, p, 0.3, 1.4, width=9.4)

    add_bullet_list(slide, 0.6, 4.3, 8.8, 1.0, [
        "Sowbot data is much simpler (single IR camera, controlled setup) → models converge quickly",
        "Despite high sowbot accuracy, fine-tuning on OAK-D didn't outperform from-scratch ImageNet",
    ], size=12, color=LIGHT_GRAY)


def make_predictions_slide(prs):
    """Slide 11: Random predictions — actual model outputs on test images."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    slide_title(slide, "Sample Predictions — Test Set",
                "Random depth images from held-out pigs with model predictions")

    p = img_path("random_preds_posture3clean_depth_test.png")
    if p:
        add_image(slide, p, 1.5, 1.3, width=7.0)


def make_colab_slide(prs):
    """Slide 12: Colab workflow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    slide_title(slide, "Google Colab Training", "GPU-accelerated workflow")

    tbl = add_table(slide, 1.5, 1.5, 6.8, 1.2, 2, 3)

    for j, h in enumerate(["", "Mac (MPS)", "Colab (T4 GPU)"]):
        style_cell(tbl.cell(0, j), h, size=13, bold=True, color=WHITE, fill_color=ACCENT)
    for j, val in enumerate(["Time per epoch", "40–400 sec", "4–26 sec"]):
        c = RED if j == 1 else GREEN if j == 2 else WHITE
        style_cell(tbl.cell(1, j), val, size=13, color=c,
                   fill_color=RGBColor(0x25, 0x25, 0x3A), bold=(j > 0))

    add_text(slide, 0.6, 3.0, 8.8, 0.4, "Workflow:", size=16, color=WHITE, bold=True)

    add_bullet_list(slide, 0.6, 3.4, 8.8, 2.0, [
        "Code syncs via GitHub (git push local → git pull on Colab)",
        "Data stored on Google Drive (upload once, mount every session)",
        "Training on T4 GPU — 10x faster than local Mac",
        "All 9 models train in ~15 minutes total on Colab",
    ], size=13, color=LIGHT_GRAY)


def make_recommendation_slide(prs):
    """Slide 13: Recommendation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    slide_title(slide, "Recommendation")

    add_rect(slide, 1.0, 1.5, 7.8, 1.4, RGBColor(0x1C, 0x3D, 0x2A), GREEN)

    add_text(slide, 1.3, 1.6, 7.2, 0.5,
             "Best Model: MobileNetV2 / Depth",
             size=24, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, 1.3, 2.1, 7.2, 0.5,
             "98.4% accuracy  |  3.4M params  |  100% sitting recall  |  works in darkness",
             size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, 0.6, 3.3, 8.8, 0.4, "Why this model:", size=16, color=WHITE, bold=True)

    add_bullet_list(slide, 0.6, 3.7, 8.8, 2.0, [
        "Same accuracy as larger models (DenseNet, Xception) with 10x fewer parameters",
        "Best sitting detection: only model to correctly classify all 4 sitting test frames",
        "Depth modality works in complete darkness (important for 24/7 barn monitoring)",
        "Smallest model = fastest inference, easiest to deploy on edge devices (OAK-D)",
    ], size=13, color=LIGHT_GRAY)


def make_next_steps_slide(prs):
    """Slide 14: Next steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    slide_title(slide, "Next Steps")

    add_text(slide, 0.6, 1.4, 8.8, 0.4,
             "Task A — Remaining (optional):", size=16, color=ACCENT, bold=True)

    add_bullet_list(slide, 0.6, 1.8, 8.8, 1.5, [
        "Label more sitting frames to improve sitting F1 (currently 52 samples)",
        "Test multi-modal fusion (RGB+Depth, 6-channel input)",
        "Export best model to ONNX for real-time deployment on OAK-D",
    ], size=13, color=LIGHT_GRAY)

    add_rect(slide, 0.5, 3.4, 9.0, 0.04, ACCENT)

    add_text(slide, 0.6, 3.6, 8.8, 0.4,
             "Task B — Vulva Segmentation:", size=16, color=ACCENT, bold=True)

    add_bullet_list(slide, 0.6, 4.0, 8.8, 1.5, [
        "Segment vulva region in depth/IR images for swelling detection",
        "Build on posture pipeline — reuse data loading and training infrastructure",
        "Ready to start upon supervisor approval",
    ], size=13, color=LIGHT_GRAY)


def make_thankyou_slide(prs):
    """Slide 15: Thank you / Questions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)

    add_text(slide, 0.6, 2.0, 8.8, 0.8,
             "Thank You",
             size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(3.5), Inches(2.9), Inches(3.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_text(slide, 0.6, 3.2, 8.8, 0.5,
             "Questions?",
             size=24, color=MID_GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, 0.6, 4.2, 8.8, 0.4,
             "github.com/AkbarDevop/paal-pipeline",
             size=14, color=ACCENT, align=PP_ALIGN.CENTER, font_name="Consolas")


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    make_title_slide(prs)          # 1.  Title
    make_pipeline_slide(prs)       # 2.  Pipeline diagram
    make_posture_examples(prs)     # 3.  Posture examples (RGB photos)
    make_modality_slide(prs)       # 4.  3 modalities × 3 postures grid
    make_crop_slide(prs)           # 5.  Crop before/after
    make_data_slide(prs)           # 6.  Dataset + chart
    make_models_slide(prs)         # 7.  Model architecture
    make_results_slide(prs)        # 8.  Results table + bar chart
    make_confusion_slide(prs)      # 9.  Confusion matrices
    make_sitting_slide(prs)        # 10. Sitting analysis + F1 chart
    make_predictions_slide(prs)    # 11. Sample predictions
    make_sowbot_slide(prs)         # 12. Sowbot comparison charts
    make_sowbot_curves_slide(prs)  # 13. Sowbot training curves
    make_colab_slide(prs)          # 14. Colab workflow
    make_recommendation_slide(prs) # 15. Recommendation
    make_next_steps_slide(prs)     # 16. Next steps
    make_thankyou_slide(prs)       # 17. Thank you

    out_path = os.path.join(OUTPUT_DIR, "weekly_report_task_a.pptx")
    prs.save(out_path)
    print(f"Presentation saved: {out_path}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
